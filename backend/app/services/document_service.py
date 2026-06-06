"""Document parsing service.

Extracts plain text from PDF, DOCX, and PPTX uploads using fully
offline, open-source libraries.  The caller is responsible for
deleting the temporary file after this module returns.
"""

from __future__ import annotations

import os
import tempfile
from dataclasses import dataclass
from pathlib import Path

import structlog

logger = structlog.get_logger(__name__)

ALLOWED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".doc"}
ALLOWED_MIME_TYPES = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/msword",
}


@dataclass
class DocumentExtract:
    text: str
    title: str
    page_count: int


def _parse_pdf(path: str) -> tuple[str, int]:
    import fitz  # PyMuPDF

    try:
        doc = fitz.open(path)
    except Exception as exc:
        raise ValueError(f"Could not read PDF: {exc}") from exc
    pages = [page.get_text() for page in doc]
    doc.close()
    return "\n".join(pages), len(pages)


def _parse_docx(path: str) -> tuple[str, int]:
    import zipfile
    from docx import Document

    try:
        doc = Document(path)
    except (zipfile.BadZipFile, Exception) as exc:
        raise ValueError(f"Could not read DOCX: {exc}") from exc
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs), len(doc.paragraphs)


def _parse_pptx(path: str) -> tuple[str, int]:
    import zipfile
    from pptx import Presentation

    try:
        prs = Presentation(path)
    except (zipfile.BadZipFile, Exception) as exc:
        raise ValueError(f"Could not read PPTX: {exc}") from exc
    slides_text: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
        if parts:
            slides_text.append("\n".join(parts))
    return "\n\n".join(slides_text), len(prs.slides)


async def parse_upload(
    filename: str,
    content: bytes,
    mime_type: str | None = None,
) -> DocumentExtract:
    """Parse an in-memory upload and return extracted text.

    Writes content to a temp file (required by the parsing libraries),
    parses it, and removes the temp file before returning.
    """
    if not content:
        raise ValueError("Uploaded file is empty.")

    suffix = Path(filename).suffix.lower()
    if suffix not in ALLOWED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type '{suffix}'. "
            f"Allowed: {', '.join(sorted(ALLOWED_EXTENSIONS))}"
        )

    title = Path(filename).stem

    tmp_path: str | None = None
    try:
        with tempfile.NamedTemporaryFile(
            suffix=suffix, delete=False
        ) as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        if suffix == ".pdf":
            text, page_count = _parse_pdf(tmp_path)
        elif suffix in (".docx", ".doc"):
            text, page_count = _parse_docx(tmp_path)
        elif suffix == ".pptx":
            text, page_count = _parse_pptx(tmp_path)
        else:
            raise ValueError(f"Unhandled suffix: {suffix}")

        logger.info(
            "document.parsed",
            filename=filename,
            chars=len(text),
            pages=page_count,
        )
        return DocumentExtract(text=text.strip(), title=title, page_count=page_count)

    finally:
        if tmp_path and os.path.exists(tmp_path):
            os.unlink(tmp_path)
